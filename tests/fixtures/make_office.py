# /// script
# requires-python = ">=3.12,<3.13"
# dependencies = [
#   "python-pptx==1.0.2",
#   "openpyxl==3.1.5",
#   "python-docx==1.1.2",
#   "pillow==11.0.0",
# ]
# ///
# The interpreter is pinned, unlike everywhere else in this repo, because
# pillow==11.0.0 ships no wheel for 3.14 and falls back to a source build that
# needs libjpeg headers. Fixture bytes must be reproducible, so the generator
# pins both the libraries and the interpreter that produced the committed
# files. Nothing here runs at skill runtime; see the module docstring.
"""Generate the Office fixtures the doc-extract tests run against.

    uv run tests/fixtures/make_office.py

GENERATION-TIME ONLY. None of these libraries may enter the skill's runtime
dependency set -- the packaging promise is uv and nothing else, and anydoc is
the text engine. They are here because they are the only way to author OOXML
packages by hand without LibreOffice.

The generated files are committed, so the tests do not depend on this script
running. Regenerate only when a fixture needs to change, and check the result
in: a test that silently regenerates its own input proves nothing.

Each fixture exists for a named failure this design already tripped over:

  deck.pptx        one logo on every slide (UBIQUITY), one photo (content),
                   one chart WITH cached values -- python-pptx writes caches,
                   which is the branch openpyxl cannot produce.
  imageheavy.pptx  10 slides x 3 distinct images: SCALE_GUARD and the
                   per-slide asset-cap behaviour.
  book.xlsx        sheets [Data, Empty, Notes]. The empty sheet is the whole
                   point: anydoc emits nothing for it, so zipping sheet names
                   against table blocks positionally cites the wrong sheet.
                   Also carries a chart with NO cache and an embedded image,
                   neither of which anydoc surfaces at all.
  single.xlsx      one sheet, so anydoc emits no heading at all and the name
                   has to come from the package.
  doc.docx         headings + a body image placed twice: distinct image
                   inlines sharing one asset id, which is where docx placement
                   counts come from.
  headless.docx    no headings whatsoever -- the citation-granularity fallback.
  collide.pptx     prose byte-identical to the placeholder anydoc emits for a
                   picture. An inline description anchors to that placeholder,
                   which is ordinary text: a slide ABOUT a file called
                   image.png, or a caption repeating a picture's alt text,
                   renders the same line the picture does. One slide per
                   flavour of the collision.
"""
import pathlib, io

HERE = pathlib.Path(__file__).resolve().parent


def _png(rgb, size=(160, 160)):
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", size, rgb).save(buf, "PNG")
    buf.seek(0)
    return buf


# A logo is small and repeated; a photo is large and unique. The furniture
# filter separates them on those two facts alone, so the fixtures differ in
# exactly those two ways and nothing else.
LOGO = lambda: _png((200, 30, 40), (90, 90))       # under MIN_DIM=120
PHOTO = lambda: _png((30, 90, 200), (640, 480))
SLIVER = lambda: _png((10, 10, 10), (900, 40))     # aspect 22.5 > MAX_ASPECT


def deck():
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    blank, titled = prs.slide_layouts[6], prs.slide_layouts[5]

    s1 = prs.slides.add_slide(titled)
    s1.shapes.title.text = "Quarterly Review"
    s1.shapes.add_picture(LOGO(), Inches(0.2), Inches(0.2))

    s2 = prs.slides.add_slide(titled)
    s2.shapes.title.text = "Site Photograph"
    s2.shapes.add_picture(LOGO(), Inches(0.2), Inches(0.2))
    s2.shapes.add_picture(PHOTO(), Inches(1), Inches(1.5), width=Inches(4))
    s2.notes_slide.notes_text_frame.text = "Photo taken during the site visit."

    s3 = prs.slides.add_slide(titled)
    s3.shapes.title.text = "Revenue and Cost"
    s3.shapes.add_picture(LOGO(), Inches(0.2), Inches(0.2))
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Revenue", (1240, 1310, 1180))
    data.add_series("Cost", (890, 902, 951))
    s3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(1), Inches(1.5), Inches(6), Inches(4), data)
    prs.save(HERE / "deck.pptx")


def collide():
    """Three ways prose can look exactly like an image placeholder.

    anydoc renders a picture as its alt text on a line of its own: the generic
    `image.png` when the author set none, the author's words when they did.
    Nothing distinguishes that line from a paragraph a human typed.

      s01  a text box reading `image.png`, ABOVE the picture -- naive
           "first placeholder on the slide" anchoring lands on the prose.
      s02  the same string inside a table cell, which is NOT a line of its
           own; whole-line matching must ignore it and still anchor the
           picture correctly.
      s03  two pictures sharing the author alt text "Site plan", plus a
           caption paragraph repeating it: three candidate lines, two
           pictures.

    Pictures are distinct colours and one slide each, so the furniture filter
    keeps all of them and the deck really does produce four pending items.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    titled = prs.slide_layouts[5]

    s1 = prs.slides.add_slide(titled)
    s1.shapes.title.text = "Naming convention"
    s1.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6),
                          Inches(0.6)).text_frame.text = "image.png"
    s1.shapes.add_picture(_png((10, 120, 60), (640, 480)), Inches(0.5), Inches(2.4),
                          width=Inches(4))

    s2 = prs.slides.add_slide(titled)
    s2.shapes.title.text = "File table"
    tbl = s2.shapes.add_table(2, 2, Inches(0.5), Inches(1.4),
                              Inches(5), Inches(1.2)).table
    tbl.cell(0, 0).text = "asset"
    tbl.cell(0, 1).text = "image.png"
    tbl.cell(1, 0).text = "size"
    tbl.cell(1, 1).text = "640x480"
    s2.shapes.add_picture(_png((200, 60, 10), (640, 480)), Inches(0.5), Inches(3.2),
                          width=Inches(4))

    s3 = prs.slides.add_slide(titled)
    s3.shapes.title.text = "Site plan"
    for n, rgb in enumerate(((40, 40, 200), (40, 200, 200))):
        pic = s3.shapes.add_picture(_png(rgb, (640, 480)),
                                    Inches(0.5 + 3.5 * n), Inches(1.6), width=Inches(3))
        pic._element._nvXxPr.cNvPr.set("descr", "Site plan")
    s3.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(6),
                          Inches(0.6)).text_frame.text = "Site plan"

    prs.save(HERE / "collide.pptx")


def imageheavy():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    titled = prs.slide_layouts[5]
    for i in range(10):
        s = prs.slides.add_slide(titled)
        s.shapes.title.text = f"Panel {i + 1}"
        for j in range(3):
            # Distinct colours so the byte-hash dedup cannot collapse them:
            # this fixture must produce 30 genuinely separate items.
            s.shapes.add_picture(_png((20 * i, 80 + 5 * j, 200 - 10 * j), (400, 300)),
                                 Inches(0.5 + 3 * j), Inches(2), width=Inches(2.5))
    prs.save(HERE / "imageheavy.pptx")


def book():
    import openpyxl
    from openpyxl.chart import BarChart, Reference
    from openpyxl.drawing.image import Image as XLImage

    wb = openpyxl.Workbook()
    d = wb.active
    d.title = "Data"
    for row in (("region", "revenue", "cost"),
                ("North", 1240, 890), ("South", 1310, 902), ("East", 1180, 951)):
        d.append(row)

    # openpyxl writes c:f range references and NO numCache, which is exactly
    # the branch the extractor must cover by resolving ranges against sheet
    # data. Excel-authored files carry caches; this one deliberately does not.
    ch = BarChart()
    ch.title = "Revenue by region"
    ch.add_data(Reference(d, min_col=2, min_row=1, max_row=4), titles_from_data=True)
    ch.set_categories(Reference(d, min_col=1, min_row=2, max_row=4))
    d.add_chart(ch, "F2")

    p = HERE / "_logo.png"
    p.write_bytes(LOGO().getvalue())
    d.add_image(XLImage(str(p)), "F20")

    wb.create_sheet("Empty")                       # emits nothing -- the point
    n = wb.create_sheet("Notes")
    n.append(("metric", "value"))
    n.append(("margin", 0.31))

    wb.save(HERE / "book.xlsx")
    p.unlink()


def single():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(("item", "eur"))
    ws.append(("travel", 2746))
    wb.save(HERE / "single.xlsx")


def docs():
    import docx
    d = docx.Document()
    d.add_heading("Introduction", level=1)
    d.add_paragraph("Opening paragraph of the report.")
    d.add_heading("Budget assumptions", level=2)
    d.add_paragraph("The assumptions below drive the model.")
    p = HERE / "_shot.png"
    p.write_bytes(PHOTO().getvalue())
    d.add_picture(str(p))
    d.add_heading("Results", level=1)
    d.add_paragraph("Closing paragraph.")
    d.add_picture(str(p))          # same bytes twice: two inlines, one asset
    d.save(HERE / "doc.docx")

    h = docx.Document()
    for t in ("A contract structured as prose, with no Heading styles at all.",
              "Clause 1. The parties agree to the following terms.",
              "Clause 2. Payment falls due within thirty days."):
        h.add_paragraph(t)
    h.add_picture(str(p))
    h.save(HERE / "headless.docx")
    p.unlink()


if __name__ == "__main__":
    deck(); imageheavy(); book(); single(); docs(); collide()
    for f in sorted(HERE.glob("*.pptx")) + sorted(HERE.glob("*.xlsx")) + sorted(HERE.glob("*.docx")):
        print(f"{f.name:<20} {f.stat().st_size:>7} bytes")
